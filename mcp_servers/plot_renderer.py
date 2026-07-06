from mcp.server.fastmcp import FastMCP
import os, json
from pathlib import Path
import matplotlib
matplotlib.use('Agg') # Non-interactive backend
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime
from mcp_servers.vision_extractor import _normalize_key
import math
import difflib

mcp = FastMCP("plot-renderer")
OUTPUT_DIR = Path("output")

def safe_float(val):
    if val is None:
        return 0.0
    try:
        f_val = float(val)
        if math.isnan(f_val) or math.isinf(f_val):
            return 0.0
        return f_val
    except (ValueError, TypeError):
        return 0.0

def make_radar_chart(exams, values, class_avg_values, title, save_path):
    num_vars = len(exams)
    if num_vars < 3:
        # Fallback to line chart if not enough data points for a polygon
        make_line_chart(exams, values, class_avg_values, False, title, save_path)
        return
        
    angles = np.linspace(0, 2 * np.pi, num_vars, endpoint=False).tolist()
    
    # Close the loop
    plot_values = values + [values[0]]
    plot_angles = angles + [angles[0]]
    
    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    
    # Plot student values
    ax.plot(plot_angles, plot_values, color='#1a73e8', linewidth=2.5, label='Student', marker='o')
    ax.fill(plot_angles, plot_values, color='#1a73e8', alpha=0.15)
    
    # Plot class averages if available
    if class_avg_values and len(class_avg_values) == num_vars:
        plot_avg = class_avg_values + [class_avg_values[0]]
        ax.plot(plot_angles, plot_avg, color='#ea4335', linewidth=2.0, linestyle='--', label='Class Average', marker='s')
        ax.fill(plot_angles, plot_avg, color='#ea4335', alpha=0.1)
        
    # Rotate plot so first variable is at the top
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    
    # Labels
    ax.set_thetagrids(np.degrees(angles), exams)
    
    # Style grid
    ax.grid(color='#f1f3f4', linestyle='-', linewidth=1)
    ax.spines['polar'].set_color('#dadce0')
    
    plt.title(title, size=14, color='#202124', y=1.1, weight='bold')
    plt.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), frameon=True, facecolor='white', edgecolor='none')
    plt.tight_layout()
    plt.savefig(save_path, dpi=150)
    plt.close(fig)

def make_line_chart(exams, values, class_avg_values, trend_line_enabled, title, save_path):
    fig, ax = plt.subplots(figsize=(8, 5))
    
    # Student scores
    ax.plot(exams, values, marker='o', color='#1a73e8', linewidth=2.5, label='Student', markersize=8)
    
    # Class average scores
    if class_avg_values and len(class_avg_values) == len(exams):
        ax.plot(exams, class_avg_values, marker='s', linestyle='--', color='#ea4335', linewidth=2, label='Class Average', markersize=6)
        
    # Trend line
    if trend_line_enabled and len(values) >= 2:
        try:
            x_indices = np.arange(len(exams))
            z = np.polyfit(x_indices, values, 1)
            p = np.poly1d(z)
            ax.plot(exams, p(x_indices), linestyle=':', color='#34a853', linewidth=2, label='Trend')
        except Exception as e:
            print(f"Warning: Failed to calculate trend line: {e}")
            
    # Styling
    ax.set_title(title, fontsize=14, pad=15, weight='bold', color='#202124')
    ax.set_xlabel('Exam / Assignment', fontsize=11, labelpad=10, color='#5f6368')
    ax.set_ylabel('Score', fontsize=11, labelpad=10, color='#5f6368')
    
    all_vals = values + (class_avg_values if class_avg_values else [])
    max_val = max(all_vals) if all_vals else 100
    min_val = min(all_vals) if all_vals else 0
    ax.set_ylim(max(0, min_val - 10), min(100 if max_val <= 100 else max_val * 1.1, max_val + 10))
    
    ax.grid(axis='y', color='#f1f3f4', linestyle='-', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#dadce0')
    ax.spines['bottom'].set_color('#dadce0')
    ax.tick_params(colors='#5f6368')
    
    ax.legend(loc='best', frameon=True, facecolor='white', edgecolor='none')
    plt.tight_layout()
    plt.savefig(save_path, dpi=150)
    plt.close(fig)

def make_bar_chart(exams, values, class_avg_values, trend_line_enabled, title, save_path):
    fig, ax = plt.subplots(figsize=(8, 5))
    
    x = np.arange(len(exams))
    
    if class_avg_values and len(class_avg_values) == len(exams):
        width = 0.35
        ax.bar(x - width/2, values, width, label='Student', color='#1a73e8')
        ax.bar(x + width/2, class_avg_values, width, label='Class Average', color='#ea4335')
        ax.set_xticks(x)
        ax.set_xticklabels(exams)
    else:
        ax.bar(exams, values, width=0.4, label='Student', color='#1a73e8')
        
    # Trend line overlay
    if trend_line_enabled and len(values) >= 2:
        try:
            x_indices = np.arange(len(exams))
            z = np.polyfit(x_indices, values, 1)
            p = np.poly1d(z)
            ax.plot(x_indices, p(x_indices), linestyle=':', color='#34a853', linewidth=2, label='Trend')
        except Exception as e:
            print(f"Warning: Failed to calculate trend line: {e}")
            
    # Styling
    ax.set_title(title, fontsize=14, pad=15, weight='bold', color='#202124')
    ax.set_xlabel('Exam / Assignment', fontsize=11, labelpad=10, color='#5f6368')
    ax.set_ylabel('Score', fontsize=11, labelpad=10, color='#5f6368')
    
    ax.grid(axis='y', color='#f1f3f4', linestyle='-', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#dadce0')
    ax.spines['bottom'].set_color('#dadce0')
    ax.tick_params(colors='#5f6368')
    
    ax.legend(loc='best', frameon=True, facecolor='white', edgecolor='none')
    plt.tight_layout()
    plt.savefig(save_path, dpi=150)
    plt.close(fig)

def make_overall_summary_chart(exams, datasets, title, save_path):
    fig, ax = plt.subplots(figsize=(10, 6))
    
    for label, values in datasets.items():
        # Clean label name
        clean_label = label.replace('_', ' ').title()
        ax.plot(exams, values, marker='o', linewidth=2, label=clean_label)
        
    ax.set_title(title, fontsize=14, pad=15, weight='bold', color='#202124')
    ax.set_xlabel('Exam / Assignment', fontsize=11, labelpad=10, color='#5f6368')
    ax.set_ylabel('Score', fontsize=11, labelpad=10, color='#5f6368')
    
    ax.grid(axis='y', color='#f1f3f4', linestyle='-', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#dadce0')
    ax.spines['bottom'].set_color('#dadce0')
    ax.tick_params(colors='#5f6368')
    
    ax.legend(loc='best', frameon=True, facecolor='white', edgecolor='none')
    plt.tight_layout()
    plt.savefig(save_path, dpi=150)
    plt.close(fig)

def get_data_for_field(data, dt_name):
    exams = []
    values = []
    class_avg_values = []
    
    norm_query = _normalize_key(dt_name)
    
    def natural_sort_key(s):
        import re
        return [int(text) if text.isdigit() else text.lower() for text in re.split(r'(\d+)', str(s.get("test_name", s.get("exam_name", ""))))]

    sorted_results = sorted(data.get("results", []), key=natural_sort_key)
    
    for r in sorted_results:
        exams.append(r.get("test_name") or r.get("exam_name", "Unknown"))
        fields = r.get("numerical_fields", {})
        avgs = r.get("class_averages", {})
        
        # Student score extraction
        val = fields.get(dt_name)
        if val is None:
            val = fields.get(norm_query)
        if val is None:
            # Exact normalized match
            for k, v in fields.items():
                if _normalize_key(k) == norm_query:
                    val = v
                    break
        if val is None:
            # Fuzzy match
            normalized_keys = { _normalize_key(k): k for k in fields.keys() }
            matches = difflib.get_close_matches(norm_query, normalized_keys.keys(), n=1, cutoff=0.8)
            if matches:
                val = fields[normalized_keys[matches[0]]]
                
        values.append(safe_float(val))
        
        # Class average extraction
        avg_val = avgs.get(dt_name)
        if avg_val is None:
            avg_val = avgs.get(norm_query)
        if avg_val is None:
            for k, v in avgs.items():
                if _normalize_key(k) == norm_query:
                    avg_val = v
                    break
        if avg_val is None:
            normalized_avg_keys = { _normalize_key(k): k for k in avgs.keys() }
            matches_avg = difflib.get_close_matches(norm_query, normalized_avg_keys.keys(), n=1, cutoff=0.8)
            if matches_avg:
                avg_val = avgs[normalized_avg_keys[matches_avg[0]]]
                
        class_avg_values.append(safe_float(avg_val))
        
    return exams, values, class_avg_values

@mcp.tool()
def render_matplotlib(data: dict, instruction_path: str) -> dict:
    """
    Renders charts for student scores based on aggregated data and user instructions.
    Saves outputs as PNG files in the output directory.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    # Read user instruction configuration
    trend_line = True
    class_avg_comparison = False
    overall_report = True
    data_types = []
    
    if os.path.exists(instruction_path):
        try:
            with open(instruction_path, 'r') as f:
                instructions = json.load(f)
            output_cfg = instructions.get("output", {})
            trend_line = output_cfg.get("trend_line", True)
            class_avg_comparison = output_cfg.get("class_avg_comparison", False)
            overall_report = output_cfg.get("overall_report", True)
            data_types = instructions.get("data_types", [])
        except Exception as e:
            print(f"Warning: Failed to load instructions from {instruction_path}: {e}")
            
    # Default to all numeric fields if data_types config is empty
    if not data_types and data.get("results"):
        all_keys = set()
        for r in data.get("results", []):
            all_keys.update(r.get("numerical_fields", {}).keys())
        data_types = [{"name": k, "enabled": True, "chart": "line"} for k in all_keys]
        
    saved_paths = []
    summary_datasets = {}
    exams = []
    
    for dt in data_types:
        if not dt.get("enabled", True):
            continue
            
        dt_name = dt.get("name")
        chart_type = dt.get("chart", "line")
        title = f"{dt_name.replace('_', ' ').title()} Progression"
        save_path = OUTPUT_DIR / f"{dt_name}.png"
        
        exams_list, values, class_avg_values = get_data_for_field(data, dt_name)
        if not exams_list:
            continue
            
        if all(v == 0.0 for v in values) and (not class_avg_comparison or all(av == 0.0 for av in class_avg_values)):
            print(f"Skipping chart for {dt_name} due to continuous zero values.")
            continue
            
        exams = exams_list # cache exams axis
        
        # Add to summary if it's a typical subject score (e.g. out of 100) and not class rank/total mark
        if dt_name not in ('total_mark', 'class_rank') and max(values) <= 120:
            summary_datasets[dt_name] = values
            
        # Draw chart depending on type
        if chart_type == "radar":
            make_radar_chart(exams, values, class_avg_values if class_avg_comparison else None, title, str(save_path))
        elif chart_type == "bar":
            make_bar_chart(exams, values, class_avg_values if class_avg_comparison else None, trend_line, title, str(save_path))
        else:
            make_line_chart(exams, values, class_avg_values if class_avg_comparison else None, trend_line, title, str(save_path))
            
        saved_paths.append(str(save_path))
        
    # Draw overall summary if enabled
    if overall_report and summary_datasets and exams:
        summary_path = OUTPUT_DIR / "overall_summary.png"
        make_overall_summary_chart(exams, summary_datasets, "Overall Subject Performance", str(summary_path))
        saved_paths.append(str(summary_path))
        
    return {
        "status": "success",
        "charts": saved_paths
    }

@mcp.tool()
def render_pptx(data: dict, instruction_path: str) -> dict:
    """
    Generates a PowerPoint presentation using python-pptx.
    Inserts student details, generated charts, and analytical bullet points.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    student_info = data.get("student", {})
    student_id = student_info.get("id", "UnknownID")
    today_str = datetime.date.today().strftime("%Y-%m-%d")
    pptx_filename = f"report_{student_id}_{today_str}.pptx"
    pptx_path = OUTPUT_DIR / pptx_filename
    
    # Read instructions
    trend_line = True
    class_avg_comparison = False
    overall_report = True
    data_types = []
    extra_description = ""
    
    if os.path.exists(instruction_path):
        try:
            with open(instruction_path, 'r') as f:
                instructions = json.load(f)
            output_cfg = instructions.get("output", {})
            trend_line = output_cfg.get("trend_line", True)
            class_avg_comparison = output_cfg.get("class_avg_comparison", False)
            overall_report = output_cfg.get("overall_report", True)
            data_types = instructions.get("data_types", [])
            extra_description = instructions.get("extra_description", "")
        except Exception as e:
            print(f"Warning: Failed to load instructions: {e}")
            
    if not data_types and data.get("results"):
        all_keys = set()
        for r in data.get("results", []):
            all_keys.update(r.get("numerical_fields", {}).keys())
        data_types = [{"name": k, "enabled": True, "chart": "line"} for k in all_keys]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Determine the assignment group display name
    assignment_group = "Student Performance"
    if data.get("results"):
        first_file = data["results"][0].get("source_file")
        if first_file:
            from agents.local_mem import detect_assignment_test, format_assignment_display_name
            assignment_test = detect_assignment_test(first_file)
            assignment_group = format_assignment_display_name(assignment_test)
        else:
            exam_name = data["results"][0].get("exam_name")
            if exam_name:
                from agents.local_mem import detect_assignment_test, format_assignment_display_name
                assignment_test = detect_assignment_test("dummy_file.pdf", exam_name)
                assignment_group = format_assignment_display_name(assignment_test)
    
    # Slide 1: Title Slide (Blank layout to have absolute control over visual elements)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Solid dark navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 39, 97) # Navy blue (#1E2761)
    
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = f"{assignment_group} Analysis"
    p_title.font.name = "Georgia"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255) # White
    p_title.space_after = Pt(24)
    
    p_details = tf.add_paragraph()
    details_text = f"Student: {student_info.get('name', 'N/A')}\n"
    if student_info.get('id'):
        details_text += f"Student ID: {student_info.get('id')}\n"
    if student_info.get('class') or student_info.get('section'):
        details_text += f"Class: {student_info.get('class', '')} {student_info.get('section', '')}\n"
    details_text += f"Report Date: {datetime.date.today().strftime('%B %d, %Y')}"
    p_details.text = details_text
    p_details.font.name = "Calibri"
    p_details.font.size = Pt(16)
    p_details.font.color.rgb = RGBColor(202, 220, 252) # Ice blue (#CADCFC)
    p_details.line_spacing = 1.3
    
    # Loop over enabled data types and create a slide for each
    for dt in data_types:
        if not dt.get("enabled", True):
            continue
            
        dt_name = dt.get("name")
        img_path = OUTPUT_DIR / f"{dt_name}.png"
        
        exams, values, class_avg_values = get_data_for_field(data, dt_name)
        if not exams:
            continue
            
        if all(v == 0.0 for v in values) and (not class_avg_comparison or all(av == 0.0 for av in class_avg_values)):
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Soft gray background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 246, 249)
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = f"{dt_name.replace('_', ' ').title()} Analysis"
        p_title.font.name = "Georgia"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(30, 39, 97)
        
        # Add white container card for the chart
        chart_card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.4), Inches(5.8), Inches(5.2)
        )
        chart_card.fill.solid()
        chart_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        chart_card.line.color.rgb = RGBColor(226, 232, 240)
        chart_card.line.width = Pt(1)
        
        # Add Image inside the card (centered/padded)
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(0.9), Inches(1.75), width=Inches(5.5), height=Inches(4.5))
            
        # Add white container card for the text analysis
        text_card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.783), Inches(1.4), Inches(5.8), Inches(5.2)
        )
        text_card.fill.solid()
        text_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        text_card.line.color.rgb = RGBColor(226, 232, 240)
        text_card.line.width = Pt(1)
        
        # Add Text Analysis box inside the card
        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.366), Inches(4.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Key Data Points:"
        p.font.name = "Georgia"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 39, 97)
        p.space_after = Pt(12)
        
        if values:
            p2 = tf.add_paragraph()
            p2.text = f"• Average Score: {np.mean(values):.2f}"
            p2.font.name = "Calibri"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(60, 64, 67)
            p2.space_after = Pt(8)
            
            p3 = tf.add_paragraph()
            p3.text = f"• Max Score: {np.max(values)} (on {exams[np.argmax(values)]})"
            p3.font.name = "Calibri"
            p3.font.size = Pt(14)
            p3.font.color.rgb = RGBColor(60, 64, 67)
            p3.space_after = Pt(8)
            
            p4 = tf.add_paragraph()
            p4.text = f"• Min Score: {np.min(values)} (on {exams[np.argmin(values)]})"
            p4.font.name = "Calibri"
            p4.font.size = Pt(14)
            p4.font.color.rgb = RGBColor(60, 64, 67)
            p4.space_after = Pt(8)
            
            if len(values) >= 2:
                diff = values[-1] - values[0]
                trend_str = "Improving" if diff > 0 else "Declining" if diff < 0 else "Stable"
                p5 = tf.add_paragraph()
                p5.text = f"• Trend Direction: {trend_str} ({diff:+.2f} points from start)"
                p5.font.name = "Calibri"
                p5.font.size = Pt(14)
                p5.font.color.rgb = RGBColor(60, 64, 67)
                p5.space_after = Pt(8)
                
            if class_avg_comparison and class_avg_values and len(class_avg_values) == len(values):
                diff_from_avg = values[-1] - class_avg_values[-1]
                comp_str = f"{abs(diff_from_avg):.2f} points above" if diff_from_avg > 0 else f"{abs(diff_from_avg):.2f} points below" if diff_from_avg < 0 else "equal to"
                p6 = tf.add_paragraph()
                p6.text = f"• Class Comparison: Student is {comp_str} the class average in the latest exam."
                p6.font.name = "Calibri"
                p6.font.size = Pt(14)
                p6.font.color.rgb = RGBColor(60, 64, 67)

    # Summary Slide
    summary_img = OUTPUT_DIR / "overall_summary.png"
    if overall_report and summary_img.exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Soft gray background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 246, 249)
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = "Overall Performance Summary"
        p_title.font.name = "Georgia"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(30, 39, 97)
        
        # Add white container card for the summary image
        chart_card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.4), Inches(5.8), Inches(5.2)
        )
        chart_card.fill.solid()
        chart_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        chart_card.line.color.rgb = RGBColor(226, 232, 240)
        chart_card.line.width = Pt(1)
        
        # Add Summary Image
        slide.shapes.add_picture(str(summary_img), Inches(0.9), Inches(1.75), width=Inches(5.5), height=Inches(4.5))
        
        # Add white container card for the analysis text
        text_card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.783), Inches(1.4), Inches(5.8), Inches(5.2)
        )
        text_card.fill.solid()
        text_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        text_card.line.color.rgb = RGBColor(226, 232, 240)
        text_card.line.width = Pt(1)
        
        # Add Text Analysis box
        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.366), Inches(4.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Key Takeaways:"
        p.font.name = "Georgia"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 39, 97)
        p.space_after = Pt(12)
        
        if extra_description:
            p2 = tf.add_paragraph()
            p2.text = f"User Notes: {extra_description}"
            p2.font.name = "Calibri"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(60, 64, 67)
            p2.space_after = Pt(12)
            
        p3 = tf.add_paragraph()
        p3.text = "Action Items:"
        p3.font.name = "Georgia"
        p3.font.bold = True
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(30, 39, 97)
        p3.space_after = Pt(8)
        
        p4 = tf.add_paragraph()
        p4.text = "1. Maintain performance in topics showing upward progression.\n2. Review test materials in subjects showing declining lines.\n3. Discuss class average discrepancies to ensure parity with peer groups."
        p4.font.name = "Calibri"
        p4.font.size = Pt(14)
        p4.font.color.rgb = RGBColor(60, 64, 67)
        p4.line_spacing = 1.3

    prs.save(str(pptx_path))
    
    return {
        "status": "success",
        "file": str(pptx_path)
    }

if __name__ == "__main__":
    mcp.run()
